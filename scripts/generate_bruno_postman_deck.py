from __future__ import annotations

from datetime import date
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BRAND_NAVY = RGBColor(16, 36, 64)
BRAND_BLUE = RGBColor(0, 102, 204)
BRAND_GRAY = RGBColor(90, 98, 110)


def _set_run(run, *, size_pt: int | None = None, bold: bool | None = None, color=None) -> None:
    if size_pt is not None:
        run.font.size = Pt(size_pt)
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    run.font.name = "Calibri"


def _set_title(slide, title: str) -> None:
    if slide.shapes.title:
        slide.shapes.title.text = title
        tf = slide.shapes.title.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        _set_run(p.runs[0], size_pt=40, bold=True, color=BRAND_NAVY)
        return
    box = slide.shapes.add_textbox(Inches(0.9), Inches(0.6), Inches(12.5), Inches(0.8))
    tf = box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    _set_run(p.runs[0], size_pt=40, bold=True, color=BRAND_NAVY)


def _add_subtitle(slide, subtitle: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.95), Inches(1.55), Inches(12.0), Inches(0.6))
    tf = box.text_frame
    tf.text = subtitle
    p = tf.paragraphs[0]
    _set_run(p.runs[0], size_pt=20, bold=False, color=BRAND_GRAY)


def _add_bullets(slide, left: float, top: float, width: float, height: float, bullets: list[str]) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = b
        p.level = 0
        p.space_after = Pt(6)
        p.alignment = PP_ALIGN.LEFT
        if p.runs:
            _set_run(p.runs[0], size_pt=22, color=BRAND_NAVY)


def _add_callout(slide, *, x: float, y: float, w: float, h: float, title: str, body: str) -> None:
    shape = slide.shapes.add_shape(
        1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE (avoid enum import for stability)
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(240, 246, 255)
    shape.line.color.rgb = RGBColor(200, 220, 245)
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p1 = tf.paragraphs[0]
    p1.text = title
    if p1.runs:
        _set_run(p1.runs[0], size_pt=18, bold=True, color=BRAND_BLUE)
    p1.space_after = Pt(6)
    p2 = tf.add_paragraph()
    p2.text = body
    if p2.runs:
        _set_run(p2.runs[0], size_pt=16, color=BRAND_NAVY)


def _add_table_slide(prs: Presentation, title: str, df: pd.DataFrame, note: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    _set_title(slide, title)

    if note:
        note_box = slide.shapes.add_textbox(Inches(0.95), Inches(1.35), Inches(12.4), Inches(0.45))
        tf = note_box.text_frame
        tf.text = note
        p = tf.paragraphs[0]
        if p.runs:
            _set_run(p.runs[0], size_pt=14, color=BRAND_GRAY)

    # Keep it readable: 5 columns max on slide
    cols = ["Challenge", "Postman", "Bruno", "Alternative", "Remarks"]
    table_df = pd.DataFrame(
        {
            "Challenge": df["Challenge"].astype(str),
            "Postman": df["Functionality in Postman"].astype(str),
            "Bruno": df["Functionality in Bruno"].astype(str),
            "Alternative": df["Alternative Option"].fillna("").astype(str),
            "Remarks": df["Remarks"].fillna("").astype(str),
        }
    )[cols]

    rows = len(table_df) + 1
    cols_n = len(cols)

    left, top, width, height = Inches(0.55), Inches(1.8), Inches(12.75), Inches(5.1)
    table = slide.shapes.add_table(rows, cols_n, left, top, width, height).table

    col_widths = [Inches(2.2), Inches(2.65), Inches(2.65), Inches(2.25), Inches(3.0)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Header style
    for c, name in enumerate(cols):
        cell = table.cell(0, c)
        cell.text = name
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND_NAVY
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                _set_run(r, size_pt=12, bold=True, color=RGBColor(255, 255, 255))

    # Body
    for r in range(len(table_df)):
        for c, col in enumerate(cols):
            val = str(table_df.iloc[r, c]).strip()
            cell = table.cell(r + 1, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if r % 2 == 0 else RGBColor(248, 250, 252)
            tf = cell.text_frame
            tf.word_wrap = True
            for p in tf.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                if p.runs:
                    _set_run(p.runs[0], size_pt=10, color=BRAND_NAVY)


def generate_deck(*, excel_path: str, out_path: str) -> None:
    df = pd.read_excel(excel_path, sheet_name=0)
    df = df.dropna(how="all")

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 widescreen
    prs.slide_height = Inches(7.5)

    today = date.today().strftime("%b %d, %Y")

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    _set_title(slide, "Postman → Bruno")
    _add_subtitle(slide, f"Tooling comparison and migration considerations • {today}")
    _add_callout(
        slide,
        x=0.95,
        y=2.55,
        w=5.9,
        h=2.25,
        title="Executive takeaway",
        body=(
            "Bruno is a local-first, Git-friendly API client that can reduce license cost and improve version control.\n"
            "Migration is straightforward for collections and CLI runs, but complex scripting, monitoring, and mock servers "
            "require alternatives."
        ),
    )
    _add_callout(
        slide,
        x=7.05,
        y=2.55,
        w=5.9,
        h=2.25,
        title="Best fit scenarios",
        body=(
            "Teams already using Git workflows\n"
            "CI-driven API checks (rather than cloud monitoring)\n"
            "Preference for lightweight tooling and repo-based collaboration"
        ),
    )

    # Slide 2: Agenda
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide, "Agenda")
    _add_bullets(
        slide,
        1.05,
        1.75,
        11.6,
        5.2,
        [
            "Why consider Bruno (and when not to)",
            "Feature comparison: Postman vs Bruno",
            "Migration approach and effort drivers",
            "Security/governance implications",
            "Recommendation and next steps",
        ],
    )

    # Slide 3: Why Bruno / Why not
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide, "Why Bruno (and why not)")
    _add_callout(
        slide,
        x=0.95,
        y=1.8,
        w=6.15,
        h=5.2,
        title="Why teams like Bruno",
        body=(
            "Local-first collections (Git-friendly)\n"
            "Simple collaboration via PRs/branching\n"
            "CLI support for CI pipelines\n"
            "Reduced reliance on cloud workspaces"
        ),
    )
    _add_callout(
        slide,
        x=7.05,
        y=1.8,
        w=6.15,
        h=5.2,
        title="Where Postman stays stronger",
        body=(
            "Richer scripting/runtime (pm.*)\n"
            "Built-in monitoring and scheduling\n"
            "Mock server + documentation generation\n"
            "Cloud governance features out-of-the-box"
        ),
    )

    # Slides 4-5: Comparison tables (split)
    df1 = df.iloc[:9].copy()
    df2 = df.iloc[9:].copy()
    _add_table_slide(prs, "Feature comparison (1/2)", df1, note="Highlights from your migration sheet")
    _add_table_slide(prs, "Feature comparison (2/2)", df2, note="Items that typically drive effort and process change")

    # Slide 6: Migration approach
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide, "Migration approach (practical)")
    _add_bullets(
        slide,
        0.95,
        1.75,
        12.3,
        5.5,
        [
            "1) Import Postman collections (JSON) into Bruno; fix any unsupported constructs.",
            "2) Convert environments to `.bru` / `.env` and align naming conventions.",
            "3) Refactor pre-request/tests: prioritize business-critical flows; simplify where possible.",
            "4) Replace monitoring with CI schedules (GitHub Actions/Azure DevOps) and standardized reports.",
            "5) Replace mock servers with a dedicated tool (WireMock/Mountebank/Mackoon) when needed.",
        ],
    )

    # Slide 7: Effort drivers (from sheet)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide, "Effort drivers (what actually takes time)")
    _add_callout(
        slide,
        x=0.95,
        y=1.8,
        w=12.25,
        h=1.25,
        title="Highest-risk migration items from your sheet",
        body="Pre-request scripts, test scripts, monitoring, mock server, secrets management",
    )
    _add_bullets(
        slide,
        1.15,
        3.2,
        12.0,
        3.9,
        [
            "Complex Postman scripting (pm.*) rarely ports 1:1 — plan manual refactoring and validation.",
            "Monitoring moves from Postman cloud to pipeline schedules; define SLAs and alerting early.",
            "Secrets must be handled explicitly (secret manager + `.env` discipline) to avoid plain-text leakage.",
            "Mocking and docs require separate tooling; align on an API contract (OpenAPI) where possible.",
        ],
    )

    # Slide 8: Recommendation & next steps
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide, "Recommendation & next steps")
    _add_callout(
        slide,
        x=0.95,
        y=1.75,
        w=12.25,
        h=1.55,
        title="Recommended path",
        body="Run a 2–3 week pilot: migrate 3–5 representative collections, wire CI runs, and validate script parity.",
    )
    _add_bullets(
        slide,
        1.15,
        3.6,
        12.0,
        3.6,
        [
            "Define success metrics: coverage migrated, CI runtime, flaky rate, developer effort, license impact.",
            "Decide on replacements: monitoring scheduler, mock server tool, reporting format, secrets management.",
            "Create a migration backlog: scripts/tests refactor, env standardization, repo layout, governance rules.",
        ],
    )

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)


def main() -> None:
    generate_deck(
        excel_path="postman_to_bruno.xlsx.xlsx",
        out_path="Bruno_Postman_Client_Deck.pptx",
    )


if __name__ == "__main__":
    main()

